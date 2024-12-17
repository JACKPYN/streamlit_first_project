import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image

# Streamlit 제목
st.title("📄 PDF to PPTX 변환기")
st.write("PDF 파일을 이미지 슬라이드로 구성된 PPTX로 변환해드립니다. 🎉")

# PDF 파일 업로드
uploaded_file = st.file_uploader("PDF 파일을 업로드하세요:", type=["pdf"])

# PDF → PPTX 변환 함수
def pdf_to_pptx(pdf_bytes):
    # PDF를 이미지로 변환
    images = convert_from_bytes(pdf_bytes)
    
    # PPTX 파일 생성
    ppt = Presentation()
    
    # 첫 번째 페이지를 기준으로 슬라이드 크기 설정
    if images:
        first_image = images[0]
        width, height = first_image.size  # 픽셀 단위 크기
        slide_width = Inches(width / 96)  # 픽셀을 인치로 변환 (1인치 = 96픽셀)
        slide_height = Inches(height / 96)
        ppt.slide_width = slide_width
        ppt.slide_height = slide_height

        # 이미지마다 슬라이드를 추가
        for i, img in enumerate(images):
            slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # 빈 슬라이드
            image_stream = BytesIO()
            img.save(image_stream, format="PNG")
            image_stream.seek(0)
            
            # 이미지 삽입
            slide.shapes.add_picture(image_stream, 0, 0, width=slide_width, height=slide_height)
    
    # PPTX 파일 메모리에 저장
    ppt_bytes = BytesIO()
    ppt.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes

# 변환 및 다운로드 버튼
if uploaded_file:
    st.success("파일 업로드 완료! 변환을 시작합니다...")
    
    # PDF를 PPTX로 변환
    pptx_file = pdf_to_pptx(uploaded_file.read())
    
    # 다운로드 버튼 제공
    st.download_button(
        label="📥 PPTX 파일 다운로드",
        data=pptx_file,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
